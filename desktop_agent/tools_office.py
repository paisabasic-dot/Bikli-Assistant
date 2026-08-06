"""
Office document creation: real Word (.docx), Excel (.xlsx), PowerPoint (.pptx).

createFile only writes plain text, so "report.docx" created that way is broken.
These tools build proper Office Open XML files the user can open in Word /
Excel / PowerPoint (or LibreOffice / Google Docs).

  createWordFile       -> .docx
  createExcelFile      -> .xlsx
  createPowerPointFile -> .pptx
"""

from __future__ import annotations

from pathlib import Path
from typing import Any, Dict, List, Optional, Sequence, Union

from .registry import ToolError, register
from .tools_files import _ensure_safe, resolve_user_path


def _require_lib(import_name: str, pip_name: str, feature: str):
    """Import an optional Office library or raise a clear ToolError."""
    try:
        return __import__(import_name)
    except ImportError as e:
        raise ToolError(
            f"Cannot create {feature}: the '{pip_name}' package is not installed. "
            f"Run: pip install {pip_name}"
        ) from e


def _prepare_path(
    path: Optional[str],
    *,
    default_suffix: str,
    overwrite: bool,
) -> Path:
    if not path:
        raise ToolError("Parameter 'path' is required.")
    p = resolve_user_path(str(path))
    if p.suffix.lower() != default_suffix.lower():
        # If user gave "report" or "report.txt", force the correct Office extension.
        p = p.with_suffix(default_suffix)
    _ensure_safe(p)
    if p.exists() and not overwrite:
        raise ToolError(f"File already exists: {p}. Pass overwrite=true to replace.")
    p.parent.mkdir(parents=True, exist_ok=True)
    return p


def _as_list(value: Any) -> List[Any]:
    if value is None:
        return []
    if isinstance(value, list):
        return value
    if isinstance(value, tuple):
        return list(value)
    return [value]


def _paragraphs_from_args(args: Dict[str, Any]) -> List[str]:
    """Build a list of body paragraphs from content / paragraphs / text."""
    paragraphs = args.get("paragraphs")
    if paragraphs is not None:
        return [str(p) for p in _as_list(paragraphs) if str(p).strip() != ""]

    content = args.get("content")
    if content is None:
        content = args.get("text")
    if content is None:
        content = args.get("body")
    if content is None:
        return []

    if isinstance(content, list):
        return [str(p) for p in content if str(p).strip() != ""]

    text = str(content).replace("\r\n", "\n").replace("\r", "\n")
    # Split on blank lines first (real paragraphs); fall back to single lines.
    blocks = [b.strip() for b in text.split("\n\n") if b.strip()]
    if len(blocks) > 1:
        return blocks
    lines = [ln.strip() for ln in text.split("\n") if ln.strip()]
    return lines if lines else ([text.strip()] if text.strip() else [])


def _rows_from_args(args: Dict[str, Any]) -> tuple[List[str], List[List[Any]]]:
    """
    Normalize spreadsheet data into (headers, rows).

    Accepted shapes:
      headers=["A","B"] + rows=[[1,2],[3,4]]
      data=[["A","B"],[1,2],[3,4]]   (first row = headers if headers not given)
      data=[{name: "x", age: 1}, ...]
      content="Name,Age\\nAlice,30\\nBob,25"  (CSV-ish)
    """
    headers: List[str] = [str(h) for h in _as_list(args.get("headers"))]
    rows_raw = args.get("rows")
    data = args.get("data")
    content = args.get("content")

    rows: List[List[Any]] = []

    if rows_raw is not None:
        for row in _as_list(rows_raw):
            if isinstance(row, dict):
                if not headers:
                    headers = [str(k) for k in row.keys()]
                rows.append([row.get(h, "") for h in headers])
            elif isinstance(row, (list, tuple)):
                rows.append(list(row))
            else:
                rows.append([row])
        return headers, rows

    if data is not None:
        data_list = _as_list(data)
        if not data_list:
            return headers, rows
        if all(isinstance(item, dict) for item in data_list):
            if not headers:
                # Preserve key order from first object, then union later keys.
                seen = []
                for item in data_list:
                    for k in item.keys():
                        sk = str(k)
                        if sk not in seen:
                            seen.append(sk)
                headers = seen
            for item in data_list:
                rows.append([item.get(h, item.get(str(h), "")) for h in headers])
            return headers, rows
        # List of lists
        matrix = []
        for item in data_list:
            if isinstance(item, (list, tuple)):
                matrix.append(list(item))
            else:
                matrix.append([item])
        if not headers and matrix:
            headers = [str(c) for c in matrix[0]]
            rows = matrix[1:]
        else:
            rows = matrix
        return headers, rows

    if content is not None:
        text = str(content).replace("\r\n", "\n").replace("\r", "\n").strip()
        if not text:
            return headers, rows
        parsed: List[List[str]] = []
        for line in text.split("\n"):
            line = line.strip()
            if not line:
                continue
            if "\t" in line:
                parsed.append([c.strip() for c in line.split("\t")])
            else:
                parsed.append([c.strip() for c in line.split(",")])
        if not headers and parsed:
            headers = [str(c) for c in parsed[0]]
            rows = parsed[1:]
        else:
            rows = parsed
        return headers, rows

    return headers, rows


def _slides_from_args(args: Dict[str, Any]) -> List[Dict[str, Any]]:
    """
    Normalize slides into a list of {title, bullets[]}.

    Accepted:
      slides=[{title, content|bullets|body}, ...]
      title=... + content=...  (single slide / title slide + body)
    """
    slides_raw = args.get("slides")
    if slides_raw is not None:
        out: List[Dict[str, Any]] = []
        for s in _as_list(slides_raw):
            if isinstance(s, str):
                out.append({"title": s, "bullets": []})
                continue
            if not isinstance(s, dict):
                out.append({"title": str(s), "bullets": []})
                continue
            title = str(s.get("title") or s.get("heading") or s.get("name") or "Slide")
            bullets = s.get("bullets") or s.get("points") or s.get("items")
            if bullets is None:
                body = s.get("content") or s.get("body") or s.get("text") or ""
                if isinstance(body, list):
                    bullets = body
                else:
                    text = str(body).replace("\r\n", "\n")
                    bullets = [ln.strip("-• \t") for ln in text.split("\n") if ln.strip()]
            else:
                bullets = [str(b) for b in _as_list(bullets)]
            out.append({"title": title, "bullets": bullets})
        return out

    # Single-slide convenience from top-level title/content.
    title = args.get("title") or args.get("heading") or "Presentation"
    body = args.get("content") or args.get("body") or args.get("text") or ""
    if isinstance(body, list):
        bullets = [str(b) for b in body]
    else:
        text = str(body).replace("\r\n", "\n")
        bullets = [ln.strip("-• \t") for ln in text.split("\n") if ln.strip()]
    return [{"title": str(title), "bullets": bullets}]


@register("createWordFile")
def create_word_file(args: Dict[str, Any]) -> Dict[str, Any]:
    """Create a real Microsoft Word .docx file."""
    docx = _require_lib("docx", "python-docx", "Word documents")
    from docx.shared import Pt  # type: ignore

    overwrite = bool(args.get("overwrite", False))
    p = _prepare_path(args.get("path"), default_suffix=".docx", overwrite=overwrite)

    Document = docx.Document
    doc = Document()

    title = args.get("title") or args.get("heading")
    if title:
        heading = doc.add_heading(str(title), level=0)
        for run in heading.runs:
            run.font.size = Pt(24)

    # Optional structured headings: [{text, level}]
    for h in _as_list(args.get("headings")):
        if isinstance(h, dict):
            text = str(h.get("text") or h.get("title") or "")
            level = int(h.get("level") or 1)
        else:
            text = str(h)
            level = 1
        if text:
            doc.add_heading(text, level=max(1, min(level, 4)))

    paragraphs = _paragraphs_from_args(args)
    if not paragraphs and not title:
        # Empty doc still useful as a starter file.
        doc.add_paragraph("")
    for para_text in paragraphs:
        # Markdown-ish: lines starting with # become headings.
        if para_text.startswith("### "):
            doc.add_heading(para_text[4:].strip(), level=3)
        elif para_text.startswith("## "):
            doc.add_heading(para_text[3:].strip(), level=2)
        elif para_text.startswith("# "):
            doc.add_heading(para_text[2:].strip(), level=1)
        elif para_text.startswith("- ") or para_text.startswith("* "):
            doc.add_paragraph(para_text[2:].strip(), style="List Bullet")
        else:
            doc.add_paragraph(para_text)

    # Optional simple table: table=[[h1,h2],[r1,r2],...]
    table_data = args.get("table")
    if table_data:
        matrix = [list(r) if isinstance(r, (list, tuple)) else [r] for r in _as_list(table_data)]
        if matrix:
            cols = max(len(r) for r in matrix)
            table = doc.add_table(rows=len(matrix), cols=cols)
            table.style = "Table Grid"
            for i, row in enumerate(matrix):
                for j in range(cols):
                    cell_val = row[j] if j < len(row) else ""
                    table.rows[i].cells[j].text = str(cell_val)

    doc.save(str(p))
    return {
        "result": f"Created Word document: {p}",
        "path": str(p),
        "folder": str(p.parent),
        "type": "docx",
    }


@register("createExcelFile")
def create_excel_file(args: Dict[str, Any]) -> Dict[str, Any]:
    """Create a real Microsoft Excel .xlsx spreadsheet."""
    openpyxl = _require_lib("openpyxl", "openpyxl", "Excel spreadsheets")
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side  # type: ignore

    overwrite = bool(args.get("overwrite", False))
    p = _prepare_path(args.get("path"), default_suffix=".xlsx", overwrite=overwrite)

    headers, rows = _rows_from_args(args)
    sheet_name = str(args.get("sheet_name") or args.get("sheet") or "Sheet1")[:31]

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = sheet_name

    header_font = Font(bold=True, color="FFFFFF")
    header_fill = PatternFill("solid", fgColor="2563EB")
    thin = Border(
        left=Side(style="thin", color="CBD5E1"),
        right=Side(style="thin", color="CBD5E1"),
        top=Side(style="thin", color="CBD5E1"),
        bottom=Side(style="thin", color="CBD5E1"),
    )

    start_row = 1
    title = args.get("title")
    if title:
        ws.cell(row=1, column=1, value=str(title))
        ws.cell(row=1, column=1).font = Font(bold=True, size=14)
        start_row = 3

    col_count = max(len(headers), max((len(r) for r in rows), default=0), 1)

    if headers:
        for col, h in enumerate(headers, start=1):
            cell = ws.cell(row=start_row, column=col, value=str(h))
            cell.font = header_font
            cell.fill = header_fill
            cell.alignment = Alignment(horizontal="center", vertical="center")
            cell.border = thin
        data_start = start_row + 1
    else:
        data_start = start_row

    for r_i, row in enumerate(rows):
        for c_i in range(col_count):
            val = row[c_i] if c_i < len(row) else ""
            # Coerce simple numeric strings to numbers for real Excel usability.
            if isinstance(val, str):
                s = val.strip()
                if s and s.replace(".", "", 1).replace("-", "", 1).isdigit():
                    try:
                        val = float(s) if "." in s else int(s)
                    except ValueError:
                        pass
            cell = ws.cell(row=data_start + r_i, column=c_i + 1, value=val)
            cell.border = thin

    # Auto-ish column widths.
    for col in range(1, col_count + 1):
        max_len = 0
        for row_cells in ws.iter_rows(min_col=col, max_col=col, min_row=start_row):
            for cell in row_cells:
                if cell.value is not None:
                    max_len = max(max_len, len(str(cell.value)))
        ws.column_dimensions[openpyxl.utils.get_column_letter(col)].width = min(max(max_len + 2, 10), 40)

    # Extra sheets if provided: sheets=[{name, headers, rows|data}, ...]
    for extra in _as_list(args.get("sheets")):
        if not isinstance(extra, dict):
            continue
        name = str(extra.get("name") or extra.get("sheet_name") or "Sheet")[:31]
        if name == ws.title:
            name = f"{name}_2"
        e_headers, e_rows = _rows_from_args(extra)
        ews = wb.create_sheet(title=name)
        if e_headers:
            for col, h in enumerate(e_headers, start=1):
                cell = ews.cell(row=1, column=col, value=str(h))
                cell.font = header_font
                cell.fill = header_fill
            base = 2
        else:
            base = 1
        for r_i, row in enumerate(e_rows):
            for c_i, val in enumerate(row, start=1):
                ews.cell(row=base + r_i, column=c_i, value=val)

    if not headers and not rows and not title:
        ws.cell(row=1, column=1, value="")

    wb.save(str(p))
    return {
        "result": f"Created Excel spreadsheet: {p}",
        "path": str(p),
        "folder": str(p.parent),
        "type": "xlsx",
        "sheet": sheet_name,
        "rows": len(rows),
        "columns": col_count if headers or rows else 0,
    }


@register("createPowerPointFile")
def create_powerpoint_file(args: Dict[str, Any]) -> Dict[str, Any]:
    """Create a real Microsoft PowerPoint .pptx presentation."""
    pptx = _require_lib("pptx", "python-pptx", "PowerPoint presentations")
    from pptx.util import Inches, Pt  # type: ignore

    overwrite = bool(args.get("overwrite", False))
    p = _prepare_path(args.get("path"), default_suffix=".pptx", overwrite=overwrite)

    Presentation = pptx.Presentation
    prs = Presentation()
    # Widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slides = _slides_from_args(args)
    presentation_title = args.get("title") or args.get("presentation_title")

    # If multiple slides were provided and a top-level title exists, make a title slide.
    if presentation_title and len(slides) >= 1 and slides[0].get("title") != str(presentation_title):
        title_layout = prs.slide_layouts[0]  # Title Slide
        slide = prs.slides.add_slide(title_layout)
        if slide.shapes.title:
            slide.shapes.title.text = str(presentation_title)
        subtitle = args.get("subtitle") or ""
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = str(subtitle)

    blank_or_bullets = prs.slide_layouts[1]  # Title and Content
    title_only = prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0]

    for s in slides:
        title = str(s.get("title") or "Slide")
        bullets: Sequence[Any] = s.get("bullets") or []
        if bullets:
            slide = prs.slides.add_slide(blank_or_bullets)
            if slide.shapes.title:
                slide.shapes.title.text = title
            # Body placeholder
            body = None
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == 1:
                    body = shape
                    break
            if body is None and len(slide.placeholders) > 1:
                body = slide.placeholders[1]
            if body is not None:
                tf = body.text_frame
                tf.clear()
                for i, bullet in enumerate(bullets):
                    text = str(bullet)
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = text
                    p.level = 0
                    p.font.size = Pt(20)
        else:
            slide = prs.slides.add_slide(title_only)
            if slide.shapes.title:
                slide.shapes.title.text = title

    if len(prs.slides) == 0:
        # Guarantee at least one slide.
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        if slide.shapes.title:
            slide.shapes.title.text = str(presentation_title or "Presentation")

    prs.save(str(p))
    return {
        "result": f"Created PowerPoint presentation: {p}",
        "path": str(p),
        "folder": str(p.parent),
        "type": "pptx",
        "slides": len(prs.slides),
    }


__all__ = [
    "create_word_file",
    "create_excel_file",
    "create_powerpoint_file",
]
